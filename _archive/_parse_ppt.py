# -*- coding: utf-8 -*-
"""解析 PPT：提取每页文本与图片数量，输出汇总 + 全文"""
import os, glob, re
from pptx import Presentation

ROOT = r'c:\Users\18315\Desktop\新建文件夹'
OUT_DIR = os.path.join(ROOT, '_ppt_text')
os.makedirs(OUT_DIR, exist_ok=True)

def extract_text(shape, out):
    if shape.shape_type == 6:  # group
        for s in shape.shapes:
            extract_text(s, out)
        return
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            t = ''.join(r.text for r in p.runs).strip()
            if t:
                out.append(t)

def parse(path):
    prs = Presentation(path)
    total_text = []
    n_img = 0
    page_summary = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            extract_text(shape, texts)
            if shape.shape_type == 13:  # picture
                n_img += 1
        page_summary.append((i, len(texts), sum(len(t) for t in texts)))
        total_text.append(f'--- Slide {i} ---')
        total_text.extend(texts)
    return total_text, n_img, page_summary

results = []
files = sorted(glob.glob(os.path.join(ROOT, '太妃ppt', '*.pptx')))
files.append(os.path.join(ROOT, '5分钟价格行为策略-南桥.pptx'))
for f in files:
    name = os.path.basename(f)
    try:
        texts, n_img, summary = parse(f)
        n_char = sum(len(t) for t in texts)
        base = re.sub(r'\.pptx$', '', name)
        with open(os.path.join(OUT_DIR, base + '.txt'), 'w', encoding='utf-8') as fh:
            fh.write('\n'.join(texts))
        results.append(f'{name}: {len(texts)} 行文本 / {n_char} 字 / {n_img} 张图片 / {len(summary)} 页')
    except Exception as e:
        results.append(f'{name}: ERROR {e}')

with open(os.path.join(ROOT, '_ppt_summary.txt'), 'w', encoding='utf-8') as fh:
    fh.write('\n'.join(results))
print('\n'.join(results))
