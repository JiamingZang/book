# -*- coding: utf-8 -*-
"""太妃PPT：提取每页图片 + RapidOCR 识别文字（图片型幻灯片 OCR）"""
import os, glob, io
from pptx import Presentation
from rapidocr_onnxruntime import RapidOCR
from PIL import Image

ROOT = r'c:\Users\18315\Desktop\新建文件夹'
PIC_DIR = os.path.join(ROOT, '_ppt_pics')
OUT_DIR = os.path.join(ROOT, '_ppt_text')
os.makedirs(PIC_DIR, exist_ok=True)

ocr = RapidOCR()


def collect_pics(shape, out):
    """递归收集图片 blob（含 group 内图片）"""
    if shape.shape_type == 6:  # group
        for s in shape.shapes:
            collect_pics(s, out)
    elif shape.shape_type == 13:  # picture
        try:
            out.append(shape.image.blob)
        except Exception:
            pass


files = sorted(glob.glob(os.path.join(ROOT, '太妃ppt', '*.pptx')))
for f in files:
    name = os.path.basename(f)
    base = os.path.splitext(name)[0]
    print(f'== {name} ==', flush=True)
    prs = Presentation(f)
    res = []
    for i, slide in enumerate(prs.slides, 1):
        blobs = []
        for shape in slide.shapes:
            collect_pics(shape, blobs)
        res.append(f'--- Slide {i} ---')
        if not blobs:
            res.append('')
            continue
        for j, blob in enumerate(blobs, 1):
            img = Image.open(io.BytesIO(blob)).convert('RGB')
            png = os.path.join(PIC_DIR, f'{base}_s{i:03d}_{j}.png')
            img.save(png)
            result, _ = ocr(png)
            txt = ''
            if result:
                txt = '\n'.join(line[1] for line in result)
            res.append(txt)
    with open(os.path.join(OUT_DIR, base + '.txt'), 'w', encoding='utf-8') as fh:
        fh.write('\n'.join(res))
    print(f'{name} done, {len(res)} lines', flush=True)
print('ALL DONE')
